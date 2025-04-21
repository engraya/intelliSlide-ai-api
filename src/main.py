from fastapi import FastAPI, BackgroundTasks
from pydantic import BaseModel
from typing import Optional
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor  # Import for color styling
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai  # Correct import
from pptx.enum.text import PP_ALIGN  # Import for text alignment
from pptx.dml.color import RGBColor  # Import for color styling
from pptx.util import Inches
from pptx.util import Inches, Pt


# Load environment variables
load_dotenv()

# Get API key for Google Gemini
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise ValueError("GOOGLE_API_KEY is not set in the environment or .env file!")

# Configure Google Gemini API
genai.configure(api_key=api_key)

app = FastAPI()

# Enable CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "https://intelli-slide-ai.vercel.app"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Request model with language support
class PPTRequest(BaseModel):
    topic: str
    num_slides: Optional[int] = 5
    language: Optional[str] = "English"  # Default language is English

# Generate slide content using Google Gemini
def generate_slide_content(topic: str, num_slides: int, language: str):
    prompt = (
        f"Generate {num_slides} slides in {language} about {topic}. "
        f"Each slide should include:\n"
        f"- A title\n"
        f"- Bullet points with short descriptions\n"
        f"- A short content or wider explanation for each bullet\n"
        f"- A placeholder for images if relevant (indicate it with 'Insert Image: [description]')\n"
        f"Format each slide as: 'Title\\n- Bullet 1 (short explanation)\\n- Bullet 2'."
    )

    # Use the Gemini model correctly
    model = genai.GenerativeModel("gemini-2.0-flash")
    response = model.generate_content(prompt)

    if response and hasattr(response, "text") and response.text:
        slides = response.text.strip().split("\n\n")  # Split into slides
        return [slide.split("\n") for slide in slides]  # Split title & bullets
    return []

def create_pptx(topic: str, slides_data, layout_preference: str = "default", filename="presentation.pptx"):
    prs = Presentation()
    
    # Define colors (adjust as needed for consistency)
    title_color = RGBColor(255, 255, 255)  # White
    content_color = RGBColor(50, 50, 50)  # Dark Gray
    bg_color = RGBColor(0, 51, 102)  # Dark Blue for contrast

    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[5]  # Title Only Layout for consistency
        slide = prs.slides.add_slide(slide_layout)

        # Adjusting Title Field
        title = slide.shapes.title
        title.text = slide_info["title"][:50]  # Ensuring title fits the slide space
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = title_color

        # Background Color (if required)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        # Adding Content Box
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(4))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, point in enumerate(slide_info["content"]):
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(24 if len(slide_info["content"]) <= 5 else 18)  # Adjust font size based on content
            p.font.color.rgb = content_color
            p.space_after = Pt(10)
            if idx == 0:
                p.font.bold = True  # First bullet point bold for emphasis
        
        # Adding Image Placeholder (if available)
        if "image_path" in slide_info and slide_info["image_path"]:
            img_path = slide_info["image_path"]
            slide.shapes.add_picture(img_path, Inches(1), Inches(4.5), width=Inches(8))  # Adjusted to fit neatly
        
    prs.save(filename)

@app.get("/")
async def welcome():
    return {"message": "Welcome to IntelliSlide-AI! Use /generate_ppt to create a PowerPoint presentation."}

@app.post("/generate_ppt")
async def generate_ppt(request: PPTRequest, background_tasks: BackgroundTasks):
    slides_data = generate_slide_content(request.topic, request.num_slides, request.language)
    filename = f"{request.topic.replace(' ', '_')}.pptx"
    background_tasks.add_task(create_pptx, request.topic, slides_data, filename)
    return {"message": "Presentation is being generated.", "filename": filename}

@app.get("/download_ppt/{filename}")
async def download_ppt(filename: str):
    file_path = f"./{filename}"
    if os.path.exists(file_path):
        return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
    return {"error": "File not found"}
